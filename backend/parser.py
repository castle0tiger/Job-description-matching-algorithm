import io
import pdfplumber
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_bytes: bytes) -> str:
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n".join(text_parts)


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    # 표 안 텍스트도 추출
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text.strip())
    return "\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            # 표 안 텍스트
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            parts.append(cell.text.strip())
    return "\n".join(parts)


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".ppt"}


def extract_text(file_bytes: bytes, filename: str) -> str:
    """파일 확장자에 따라 적절한 파서를 선택합니다."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == ".pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == ".docx":
        return extract_text_from_docx(file_bytes)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_bytes)
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext} (지원: PDF, DOCX, PPTX)")
